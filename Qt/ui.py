# -*- coding: utf-8 -*-

# Form implementation generated from reading ui file 'main.ui'
#
# Created by: PyQt5 UI code generator 5.9
#
# WARNING! All changes made in this file will be lost!

from PyQt5 import QtCore, QtGui, QtWidgets
from tkinter import Tk
from tkinter.filedialog import askopenfilename
try:
    from pptx import Presentation
except Exception as e:
    print(e)

###########CONFIG#######
PRINT_RESULTS_IN_CONSOLE = True
###########CONFIG#######



class Ui_root(object):

    def __init__(self, Pres):
        self.Pres = Pres

    def setupUi(self, root):


        root.setObjectName("root")
        root.resize(773, 540)
        root.setTabShape(QtWidgets.QTabWidget.Rounded)
        self.centralwidget = QtWidgets.QWidget(root)
        self.centralwidget.setObjectName("centralwidget")
        self.verticalLayout = QtWidgets.QVBoxLayout(self.centralwidget)
        self.verticalLayout.setObjectName("verticalLayout")
        self.verticalLayout_2 = QtWidgets.QVBoxLayout()
        self.verticalLayout_2.setObjectName("verticalLayout_2")
        self.horizontalLayout_2 = QtWidgets.QHBoxLayout()
        self.horizontalLayout_2.setObjectName("horizontalLayout_2")
        self.verticalLayout_4 = QtWidgets.QVBoxLayout()
        self.verticalLayout_4.setObjectName("verticalLayout_4")
        spacerItem = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        self.verticalLayout_4.addItem(spacerItem)
        self.WordsLabel = QtWidgets.QLabel(self.centralwidget)
        self.WordsLabel.setObjectName("WordsLabel")
        self.verticalLayout_4.addWidget(self.WordsLabel, 0, QtCore.Qt.AlignRight)
        spacerItem1 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        self.verticalLayout_4.addItem(spacerItem1)
        self.LinesLabel = QtWidgets.QLabel(self.centralwidget)
        self.LinesLabel.setObjectName("LinesLabel")
        self.verticalLayout_4.addWidget(self.LinesLabel, 0, QtCore.Qt.AlignRight)
        spacerItem2 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        self.verticalLayout_4.addItem(spacerItem2)
        self.CharLabelWOS = QtWidgets.QLabel(self.centralwidget)
        self.CharLabelWOS.setObjectName("CharLabelWOS")
        self.verticalLayout_4.addWidget(self.CharLabelWOS, 0, QtCore.Qt.AlignRight)
        spacerItem3 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        self.verticalLayout_4.addItem(spacerItem3)
        self.CharLabelWS = QtWidgets.QLabel(self.centralwidget)
        self.CharLabelWS.setObjectName("CharLabelWS")
        self.verticalLayout_4.addWidget(self.CharLabelWS, 0, QtCore.Qt.AlignRight)
        spacerItem4 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        self.verticalLayout_4.addItem(spacerItem4)
        self.horizontalLayout_2.addLayout(self.verticalLayout_4)
        self.verticalLayout_3 = QtWidgets.QVBoxLayout()
        self.verticalLayout_3.setObjectName("verticalLayout_3")
        spacerItem5 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        self.verticalLayout_3.addItem(spacerItem5)
        self.WordsValueLabel = QtWidgets.QLabel(self.centralwidget)
        self.WordsValueLabel.setText("")
        self.WordsValueLabel.setObjectName("WordsValueLabel")
        self.verticalLayout_3.addWidget(self.WordsValueLabel)
        spacerItem6 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        self.verticalLayout_3.addItem(spacerItem6)
        self.LinesValueLabel = QtWidgets.QLabel(self.centralwidget)
        self.LinesValueLabel.setText("")
        self.LinesValueLabel.setObjectName("LinesValueLabel")
        self.verticalLayout_3.addWidget(self.LinesValueLabel)
        spacerItem7 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        self.verticalLayout_3.addItem(spacerItem7)
        self.CharWOSValueLabel = QtWidgets.QLabel(self.centralwidget)
        self.CharWOSValueLabel.setText("")
        self.CharWOSValueLabel.setObjectName("CharWOSValueLabel")
        self.verticalLayout_3.addWidget(self.CharWOSValueLabel)
        spacerItem8 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        self.verticalLayout_3.addItem(spacerItem8)
        self.CharWSValueLabel = QtWidgets.QLabel(self.centralwidget)
        self.CharWSValueLabel.setText("")
        self.CharWSValueLabel.setObjectName("CharWSValueLabel")
        self.verticalLayout_3.addWidget(self.CharWSValueLabel)
        spacerItem9 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        self.verticalLayout_3.addItem(spacerItem9)
        self.horizontalLayout_2.addLayout(self.verticalLayout_3)
        self.verticalLayout_2.addLayout(self.horizontalLayout_2)
        self.horizontalLayout = QtWidgets.QHBoxLayout()
        self.horizontalLayout.setObjectName("horizontalLayout")
        self.btnOpen = QtWidgets.QPushButton(self.centralwidget)
        self.btnOpen.setObjectName("btnOpen")
        self.horizontalLayout.addWidget(self.btnOpen)
        spacerItem10 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        self.horizontalLayout.addItem(spacerItem10)
        self.btnCount = QtWidgets.QPushButton(self.centralwidget)
        self.btnCount.setMinimumSize(QtCore.QSize(0, 25))
        self.btnCount.setObjectName("btnCount")
        self.horizontalLayout.addWidget(self.btnCount)
        self.verticalLayout_2.addLayout(self.horizontalLayout)
        spacerItem11 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        self.verticalLayout_2.addItem(spacerItem11)
        self.logger = QtWidgets.QTextBrowser(self.centralwidget)
        self.logger.setObjectName("logger")
        self.verticalLayout_2.addWidget(self.logger)
        self.verticalLayout.addLayout(self.verticalLayout_2)
        root.setCentralWidget(self.centralwidget)
        self.statusbar = QtWidgets.QStatusBar(root)
        self.statusbar.setObjectName("statusbar")
        root.setStatusBar(self.statusbar)
        self.btnOpen.clicked.connect(self.openFile)
        self.btnCount.clicked.connect(self.countFile)

        self.retranslateUi(root)
        QtCore.QMetaObject.connectSlotsByName(root)

    def openFile(self):

        #get file chooser
        Tk().withdraw()
        filename = askopenfilename()

        if not(filename  == "" or filename  == "()" or filename  == None):
            self.logger.append("Trying to load: " + str(filename) + "...")

            #try to load presentation
            try:
                Pres = Presentation(filename)
                self.logger.append(filename + " loaded successfully!")
                self.Pres = Pres

            #if loading fails
            except Exception as exc:
                self.logger.append("Failed to load presentation! Check console for error code.")
                print(exc)
        else:
            pass

    def countFile(self):
        if not(self.Pres  == None ):

            self.logger.append("counting...")

            p = self.Pres

            #list for all strings
            text = []

            #get all text
            for slide in p.slides: #each slide
                for shape in slide.shapes:  #each shape
                    if(not shape.has_text_frame): #continue if the shape has no text frames
                        continue
                    for paragraph in shape.text_frame.paragraphs: #for each paragraph

                        for run in paragraph.runs: #for each run
                            text.append(run.text)  #get text and store



            #create lists to store all words
            allWords = []
            allWordsTemp = []

            #create variables to count all chars (without spaces)
            CharsWhithoutSpacesTemp = []
            Charcount = 0

            #create variables to count chars with spaces
            allTextOneString = []


            #each text line
            for n in range(len(text)):

                tempLine = text[n]

                #each word in tempLine
                for i in range(len(tempLine.split(" "))):

                    #add each word
                    tempWordsSplit = tempLine.split(" ")
                    allWordsTemp.append(tempWordsSplit[i])

            #remove word if its just a space
            for n in range(len(allWordsTemp)):

                #if the word is a space it may not be counted
                if(allWordsTemp[n] == "" or allWordsTemp[n] == " "):
                        pass

                else:
                    allWords.append(allWordsTemp[n])


            #get list for all caracters without spaces
            allTextOneString = "".join(allWords)

            self.logger.append("counted!")

            #set labels
            self.CharWOSValueLabel.setText(str(len(allTextOneString)))
            self.LinesValueLabel.setText(str(len(text)))
            self.WordsValueLabel.setText(str(len(allWords)))
            self.CharWSValueLabel.setText(str(len("".join(text))))

            if(PRINT_RESULTS_IN_CONSOLE):
                #print results
                print("-----------------------------------------------------------------")
                print("Total words: " + str(len(allWords)))
                print("Total lines: " + str(len(text)))
                print("Total caracters (without spaces): " + str(len(allTextOneString)))
                print("Total caracters (with spaces): " + str(len("".join(text))))
                print("-----------------------------------------------------------------")

        else:
            self.logger.append("Please open a presentation first.")



    def retranslateUi(self, root):
        _translate = QtCore.QCoreApplication.translate
        root.setWindowTitle(_translate("root", "DocuCount"))
        self.WordsLabel.setText(_translate("root", "Total words: "))
        self.LinesLabel.setText(_translate("root", "Total lines: "))
        self.CharLabelWOS.setText(_translate("root", "Total characters (without spaces): "))
        self.CharLabelWS.setText(_translate("root", "Total characters (with spaces): "))
        self.btnOpen.setText(_translate("root", "open presentation"))
        self.btnCount.setText(_translate("root", "count"))
