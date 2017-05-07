#!/usr/bin/env python
import platform
#import python-pptx
try:
    from pptx import Presentation
except Exception as exc:

    print("python-pptx is required!")
    OS = platform.system()

    #for windows
    if(OS == "Windows"):
        print("1. Open CMD (search for cmd and press enter).")
        print("2. Enter \"py -m pip install python pptx\".")
        print("If this does not work enter \"py -m easy_install python-pptx\".")
        exit(1)

    #for mac
    elif(OS == "Darwin"):
        print("1. Run \"python -m pip install python-pptx\" in a terminal.")
        print("2. As he asks you to install \"Command Line developer tools\", hit install (About 150MB) .")
        print("2. If necessary enter password")
        print("3. Confirm all questions showing up in the terminal. Usually you only have to type \"Y\" and hit enter.")
        print("4. After the first start of DocuCount you can delete the command line tools (Delete the folder at Library\Developer\CommandLineTools). The directory might have a slightly differing name")
        exit(1)

    #for linux
    elif(OS == "Linux"):
        print("1. Open a terminal.")
        print("2. Run \"sudo pip install python-pptx\".")
        print("3. If necessary enter password.")
        print("4. Confirm all questions showing up in the terminal. Usually you only have to type \"Y\" and hit enter.")
        exit(1)




from tkinter import Tk
from tkinter.filedialog import askopenfilename



def count(p):

    #list for all strings
    text = []

    #get all text
    for slide in p.slides: #each slide
        for shape in slide.shapes:  #each shape
            if(not shape.has_text_frame): #continue if the shape has no text frames
                continue
            for paragraph in shape.text_frame.paragraphs: #for each paragraph

                for run in paragraph.runs: #for each run
                    text.append(run.text)  #get text and store

    return text


#####MAIN#####



while True:

    pres = "null"
    while(pres == "null"):

        #get file chooser
        Tk().withdraw()
        filename = askopenfilename()

        if(filename == ""):
            break

        try:

            print("reading file...")

            #get presentation
            pres = Presentation(filename)

            #count file
            text = count(pres)
            print("counting presentation please wait...")


            #create lists to store all words
            allWords = []
            allWordsTemp = []

            #create variables to count all chars (without spaces)
            CharsWhithoutSpacesTemp = []
            Charcount = 0

            #create variables to count chars with spaces
            allTextOneString = []


            #each text line
            for n in range(len(text)):

                tempLine = text[n]

                #each word in tempLine
                for i in range(len(tempLine.split(" "))):

                    #add each word
                    tempWordsSplit = tempLine.split(" ")
                    allWordsTemp.append(tempWordsSplit[i])

            #remove word if its just a space
            for n in range(len(allWordsTemp)):

                #if the word is a space it may not be counted
                if(allWordsTemp[n] == "" or allWordsTemp[n] == " "):
                        pass

                else:
                    allWords.append(allWordsTemp[n])


            #get list for all caracters without spaces
            allTextOneString = "".join(allWords)


            #print results
            print("-----------------------------------------------------------------")
            print("Total words: " + str(len(allWords)))
            print("Total lines: " + str(len(text)))
            print("Total caracters (without spaces): " + str(len(allTextOneString)))
            print("Total caracters (with spaces): " + str(len("".join(text))))
            print("-----------------------------------------------------------------")

            #exit
            cmd = input("Hit enter to count another document, enter \"q\" to exit >>> ")
            if(cmd == "q"):
                print("Thank you for using DocuCount!")
                exit(0)


        except Exception as exc:

            print("Error! Did you chose a .pptx file?")
            pres = "null"
            print("Error Code: " + exc)
